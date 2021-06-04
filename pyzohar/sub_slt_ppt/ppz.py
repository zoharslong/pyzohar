#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Created on Mon May 31 10:04:01 2021
ppt generating operation.
requirements = [
  'python-pptx>=0.6.19',
]
@author: zoharslong
"""
from pyzohar.sub_slt_bsc.bsz import lsz
from pyzohar.sub_slt_bsc.dfz import dfz
from pandas import DataFrame, isnull
from pptx import Presentation
from os.path import join as os_join, exists as os_exists
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt, Inches
from pptx.presentation import Presentation as typ_presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_LEGEND_POSITION, XL_LABEL_POSITION
dct_plt = {
    'blk': (0, 0, 0),  # 黑
    'wht': (255, 255, 255),  # 白
    'gry': (191, 191, 191),  # 灰
    'gry_dep': (127, 127, 127),  # 深灰
    'red_lgt': (225, 178, 172),  # 浅红
    'red_nrm': (209, 136, 126),  # 红
    'red_dep': (167, 91, 81),  # 深红
}  # 调色板


class ppz_nit(dfz):
    """
    self.lcn = {fld, fls, ppt:{mtr,sld,...}}
    """
    dct_util = {
        'cm': Cm,
        'pt': Pt,
        'nch': Inches
    }              # 尺度单位
    dct_anchor = {
        'bottom': MSO_ANCHOR.BOTTOM,
        'middle': MSO_ANCHOR.MIDDLE,
        'center': MSO_ANCHOR.MIDDLE,
        'top': MSO_ANCHOR.TOP
    }            # 纵向对齐
    dct_align = {
        'justify': PP_ALIGN.JUSTIFY,  # 两端
        'left': PP_ALIGN.LEFT,
        'right': PP_ALIGN.RIGHT,
        'center': PP_ALIGN.CENTER,  # 居中
        'middle': PP_ALIGN.CENTER,  # 居中
    }             # 横向对齐
    dct_legend_position = {
        'rightout': [XL_LEGEND_POSITION.RIGHT, False],
        'rightin': [XL_LEGEND_POSITION.RIGHT, True],
    }   # 图例位置
    dct_label_position = {
        'inside_base': XL_LABEL_POSITION.INSIDE_BASE,
        'inside_end': XL_LABEL_POSITION.INSIDE_END,
        'outside_end': XL_LABEL_POSITION.OUTSIDE_END,
        'best_fit': XL_LABEL_POSITION.BEST_FIT,
        'center': XL_LABEL_POSITION.CENTER,
        'above': XL_LABEL_POSITION.ABOVE,
        'right': XL_LABEL_POSITION.RIGHT,
    }    # 标签位置
    dct_mtr = {
        'wdh': 9144000,
        'wtd': '16:10',
        'utl': 'cm',
        'fontName': 'Microsoft YaHei UI',
        'fontColor': (64, 64, 64),
        'fontSize': Pt(12),
        'fontSize_label': Pt(10),
        'fontSize_label_2': Pt(8),
        'fontBold': False,
        'fontItalic': False,
        'anchor': 'middle',  # 纵向对齐 - 居中
        'align': 'justify',  # 横向对齐 - 两端
        'line_spacing': 1.15,
        'space_before': Pt(3),
        'space_after': Pt(3),
        'mtr': {
            'full': [
                {
                    'typ': 'textbox', 'tip': '顶栏总结',
                    'left': 1.4 / 2, 'top': 0.875 / 2, 'width': 25.4 - 1.4 - 4.5, 'height': 1.5,
                    'anchor': 'middle', 'align': 'justify',
                },  # 顶栏总结 - 需要{paragraph:[{ctt}]}
                {
                    'tip': '左半内容',
                    'left': 1.4 / 2, 'top': 0.875 / 2 + 1.5 + 1, 'width': (25.4 - 1.4),
                    'height': 15.875 - (0.875 + 1.5 + 1),
                },  # 全页内容 - 需要{typ, ..}
                {
                    'typ': 'picture', 'tip': '右上图标',
                    'left': 25.4 - 4.5, 'top': 0.875 / 2, 'height': 1.5, 'width': 4.5,
                },  # 右上图标 - 需要图片地址的{fld, fls}
            ],    # 整版呈现布局[title, full, logo]
            'middle': [
                {
                    'typ': 'textbox', 'tip': '顶栏总结',
                    'left': 1.4 / 2, 'top': 0.875 / 2, 'width': 25.4 - 1.4 - 4.5, 'height': 1.5,
                    'anchor': 'middle', 'align': 'justify',
                },  # 顶栏总结 - 需要{paragraph:[{ctt}]}
                {
                    'tip': '左半内容',
                    'left': 1.4 / 2, 'top': 0.875 / 2 + 1.5 + 1, 'width': (25.4 - 1.4) / 2 - 0.25,
                    'height': 15.875 - (0.875 + 1.5 + 1),
                },  # 左半内容 - 需要{typ, ..}
                {
                    'tip': '右半内容',
                    'left': 25.4 / 2 + 0.25, 'top': 0.875 / 2 + 1.5 + 1, 'width': (25.4 - 1.4) / 2 - 0.25,
                    'height': 15.875 - (0.875 + 1.5 + 1),
                },  # 右半内容 - 需要{typ, ..}
                {
                    'typ': 'picture', 'tip': '右上图标',
                    'left': 25.4 - 4.5, 'top': 0.875 / 2, 'height': 1.5, 'width': 4.5,
                },  # 右上图标 - 需要图片地址的{fld, fls}
                {
                    'typ': 'line', 'tip': '隔断中线',
                    'left': 25.4 / 2 - 0.001, 'top': 0.875 / 2 + 1.5 + 1 + 3, 'width': 0.002,
                    'height': 15.875 - (0.875 + 1.5 + 1) - 6,
                    'color': (225, 225, 225),
                },  # 隔断中线 - 可省略
            ],  # 中位分割布局[title, left, right, logo, splitLine]
            'center': [
                {
                    'typ': 'textbox', 'tip': '顶栏总结',
                    'left': 1.4 / 2, 'top': 0.875 / 2, 'width': 25.4 - 1.4 - 4.5, 'height': 1.5,
                    'anchor': 'middle', 'align': 'justify',
                },  # 顶栏总结 - 需要{paragraph:[{ctt}]}
                {
                    'tip': '上半内容',
                    'left': 1.4 / 2, 'top': 0.875 / 2 + 1.5 + 1,
                    'width': (25.4 - 1.4),
                    'height': (15.875 - (0.875 + 1.5 + 1))/2 - 0.875/2,
                },  # 右半内容 - 需要{typ, ..}
                {
                    'tip': '下半内容',
                    'left': 1.4 / 2, 'top': 0.875 / 2 + 1.5 + 1 + (15.875 - (0.875 + 1.5 + 1))/2 + 0.875/2,
                    'width': (25.4 - 1.4),
                    'height': (15.875 - (0.875 + 1.5 + 1))/2 - 0.875/2,
                },  # 右半内容 - 需要{typ, ..}
                {
                    'typ': 'picture', 'tip': '右上图标',
                    'left': 25.4 - 4.5, 'top': 0.875 / 2, 'height': 1.5, 'width': 4.5,
                },  # 右上图标 - 需要图片地址的{fld, fls}
                {
                    'typ': 'line', 'tip': '隔断横线',
                    'left': 1.4 / 2 + 3, 'top': 0.875 / 2 + 1.5 + 1 + (15.875 - (0.875 + 1.5 + 1))/2 - 0.001,
                    'width': 25.4 - 1.4 - 6,
                    'height': 0.002,
                    'color': (225, 225, 225),
                },  # 隔断中线 - 可省略
            ],
            'bottom': [],
            'left': [
                {
                    'typ': 'textbox',
                    'tip': '顶栏总结',
                    'left': 1.4 / 2,
                    'top': 0.875 / 2,
                    'width': 25.4 - 1.4 - 4.5,
                    'height': 1.5,
                    'anchor': 'middle',
                    'align': 'justify',
                },  # 顶栏总结 - 需要{paragraph:[{ctt}]}
                {
                    'tip': '左半内容',
                    'left': 1.4 / 2,
                    'top': 0.875 / 2 + 1.5 + 1,
                    'width': (25.4 - 1.4) / 3*2 - 0.25,
                    'height': 15.875 - (0.875 + 1.5 + 1),
                },  # 左半内容 - 需要{typ, ..}
                {
                    'tip': '右半内容',
                    'left': 1.4 / 2 + (25.4-1.4) / 3*2 + 0.25,
                    'top': 0.875 / 2 + 1.5 + 1,
                    'width': (25.4 - 1.4) / 3*1 - 0.25,
                    'height': 15.875 - (0.875 + 1.5 + 1),
                },  # 右半内容 - 需要{typ, ..}
                {
                    'typ': 'picture', 'tip': '右上图标',
                    'left': 25.4 - 4.5, 'top': 0.875 / 2, 'height': 1.5, 'width': 4.5,
                },  # 右上图标 - 需要图片地址的{fld, fls}
                {
                    'typ': 'line', 'tip': '隔断中线',
                    'left': 1.4/2 + (25.4-1.4) / 3*2 - 0.001,
                    'top': 0.875 / 2 + 1.5 + 1 + 3,
                    'width': 0.002,
                    'height': 15.875 - (0.875 + 1.5 + 1) - 6,
                    'color': (225, 225, 225),
                },  # 隔断中线 - 可省略
            ],
            'right': [
                {
                    'typ': 'textbox',
                    'tip': '顶栏总结',
                    'left': 1.4 / 2,
                    'top': 0.875 / 2,
                    'width': 25.4 - 1.4 - 4.5,
                    'height': 1.5,
                    'anchor': 'middle',
                    'align': 'justify',
                },  # 顶栏总结 - 需要{paragraph:[{ctt}]}
                {
                    'tip': '左半内容',
                    'left': 1.4 / 2,
                    'top': 0.875 / 2 + 1.5 + 1,
                    'width': (25.4 - 1.4) / 3*1 - 0.25,
                    'height': 15.875 - (0.875 + 1.5 + 1),
                },  # 左半内容 - 需要{typ, ..}
                {
                    'tip': '右半内容',
                    'left': 1.4 / 2 + (25.4-1.4) / 3*1 + 0.25,
                    'top': 0.875 / 2 + 1.5 + 1,
                    'width': (25.4 - 1.4) / 3*2 - 0.25,
                    'height': 15.875 - (0.875 + 1.5 + 1),
                },  # 右半内容 - 需要{typ, ..}
                {
                    'typ': 'picture', 'tip': '右上图标',
                    'left': 25.4 - 4.5, 'top': 0.875 / 2, 'height': 1.5, 'width': 4.5,
                },  # 右上图标 - 需要图片地址的{fld, fls}
                {
                    'typ': 'line', 'tip': '隔断中线',
                    'left': 1.4/2 + (25.4-1.4) / 3*1 - 0.001,
                    'top': 0.875 / 2 + 1.5 + 1 + 3,
                    'width': 0.002,
                    'height': 15.875 - (0.875 + 1.5 + 1) - 6,
                    'color': (225, 225, 225),
                },  # 隔断中线 - 可省略
            ],
        },  # 可以直接在slide调用的布局, 通过lcn.ppt.sld.mtr指定
        'tbl': {
            'mrg': [],
            'cll_fit': {},
            'cll_fit_clm': {},
            'anchor': 'middle',
            'align': 'middle',
            'color': (225, 225, 225),
            'fontName': 'Microsoft YaHei',
            'fontColor': (0, 0, 0),
            'fontSize': 12,  # fixed Pt
            'fontBold': False,
            'fontItalic': False,
            'margin_left': 0.15,     # 跟mtr.utl
            'margin_right': 0.15,    # 跟mtr.utl
        },
    }               # 母版的默认参数
    dct_plt = dct_plt             # 调色板

    def __init__(self, dts=None, lcn=None, *, spr=False, cvr=True, ppt=None):
        super(ppz_nit, self).__init__(dts, lcn, spr=spr)
        self.__ppt = None
        self.ppt = ppt
        self.ppt_nit(cvr=cvr)

    @property
    def ppt(self):
        """
        self.ppt.
        :return: self.__ppt
        """
        return self.__ppt

    @ppt.setter
    def ppt(self, ppt):
        """
        set self.__ppt in self.ppt.
        :param ppt: main object from python-pptx
        :return: None
        """
        if type(ppt) in [typ_presentation, type(None)]:
            self.__ppt = ppt
        else:
            raise TypeError('info: ppt\'s type %s is not available.' % type(ppt))

    def ppt_nit_chk_lcn(self):
        if 'ppt' not in self.lcn.keys():
            self.lcn.update({'ppt': {'mtr': self.dct_mtr}})
        elif 'mtr' not in self.lcn['ppt'].keys():
            self.lcn['ppt'].update({'mtr': self.dct_mtr})
        else:
            dct_tmp = self.dct_mtr.copy()
            dct_tmp.update(self.lcn['ppt']['mtr'])
            self.lcn['ppt']['mtr'].update(dct_tmp)

    def ppt_nit(self, *, cvr=True):
        """
        pptx initiation,
        create a new empty self.ppt or get ppt from file
        :param cvr: create a new empty self.ppt to cover the old one in local file, default True
        """
        fls = os_join(self.lcn['fld'], self.lcn['fls']) if \
            os_exists(os_join(self.lcn['fld'], self.lcn['fls'])) and cvr is False else None
        self.ppt = Presentation(fls) if \
            self.ppt is None or type(self.ppt) not in [typ_presentation] or cvr is True else self.ppt
        self.ppt_nit_chk_lcn()

    def ppt_siz(self, wdh=None, wtd=None):
        """
        pptx's size, width and height default [10Inches, 25.4Cm, 9144000], [6.25Inches, 15.875Cm, 5715000]
        :param wdh: width
        :param wtd: width to height, 宽高比
        """
        if 'wdh' in self.lcn['ppt']['mtr'].keys():
            self.lcn['ppt']['mtr']['wdh'] = wdh = wdh if \
                wdh is not None and self.lcn['ppt']['mtr']['wdh'] is None else self.lcn['ppt']['mtr']['wdh']
        else:
            self.lcn['ppt']['mtr']['wdh'] = wdh
        if 'wtd' in self.lcn['ppt']['mtr'].keys():
            self.lcn['ppt']['mtr']['wtd'] = wtd = wtd if \
                wtd is not None and self.lcn['ppt']['mtr']['wtd'] is None else self.lcn['ppt']['mtr']['wtd']
        else:
            self.lcn['ppt']['mtr']['wtd'] = wtd
        self.ppt.slide_width = wdh
        self.ppt.slide_height = int(wdh / float(wtd.split(':')[0]) * float(wtd.split(':')[1]))

    def sld_mtr(self, ndx_sld):
        """"""
        if 'mtr' not in self.lcn['ppt'][ndx_sld] or self.lcn['ppt'][ndx_sld]['mtr'] is None:
            pass
        else:
            print('info: mother mask %s synchronization' % self.lcn['ppt'][ndx_sld]['mtr'])
            mtr = self.lcn['ppt'][ndx_sld]['mtr']
            n = 0
            for dct_i in self.dct_mtr['mtr'][mtr]:
                dct_i_tmp = dct_i.copy()
                try:
                    dct_i_tmp.update(self.lcn['ppt'][ndx_sld]['shp'][n])
                    self.lcn['ppt'][ndx_sld]['shp'][n].update(dct_i_tmp)
                except IndexError:
                    self.lcn['ppt'][ndx_sld]['shp'].append(dct_i_tmp)
                finally:
                    n += 1


class tblMixin(ppz_nit):

    def __init__(self, dts=None, lcn=None, *, spr=False, cvr=True):
        super(tblMixin, self).__init__(dts, lcn, spr=spr, cvr=cvr)

    def tbl_mtr(self, dct_shp):
        dct_tmp = self.dct_mtr['tbl'].copy()
        dct_tmp.update(dct_shp)
        dct_shp.update(dct_tmp)
        if 'row' not in dct_shp:
            dct_shp['row'], dct_shp['clm'] = DataFrame(dct_shp['cll']).shape

    def fit_mtr(self, dct_shp, tgt_fit='cll_fit', flt_len=None, *, dct_fit=None):
        """
        对dct_shp.cll_fit进行批量处理
        """
        flt_len = dct_shp['clm'] if flt_len is None else flt_len
        dct_shp[tgt_fit] = dct_fit if dct_fit is not None else dct_shp[tgt_fit]
        dct_shp[tgt_fit] = self.dct_mtr['tbl'][tgt_fit] if tgt_fit not in dct_shp else dct_shp[tgt_fit]
        for i_fit in dct_shp[tgt_fit]:
            for i_row in range(len(dct_shp[tgt_fit][i_fit])):
                if type(dct_shp[tgt_fit][i_fit][i_row]) not in [list]:
                    dct_shp[tgt_fit][i_fit][i_row] = lsz([dct_shp[tgt_fit][i_fit][i_row]]).cpy_tal(flt_len, rtn=True)

    def _it_cll(self, dct_shp, str_tgt='cll_fit'):
        """
        fit cell format operation part
        :param str_tgt: 选择是对行还是列进行批量修改, in ['cll_fit','cll_fit_clm'] for [行, 列]
        """
        dct_shp[str_tgt] = self.dct_mtr['tbl'][str_tgt] if str_tgt not in dct_shp else dct_shp[str_tgt]
        for i_fit in dct_shp[str_tgt]:
            for i_row in range(len(dct_shp[str_tgt][i_fit])):
                for i_clm in range(len(dct_shp[str_tgt][i_fit][i_row])):
                    if type(dct_shp['cll'][i_row][i_clm]) not in [dict]:
                        dct_shp['cll'][i_row][i_clm] = {'ctt': dct_shp['cll'][i_row][i_clm]}
                        dct_shp['cll'][i_row][i_clm]['ctt'] = str(dct_shp['cll'][i_row][i_clm]['ctt']) if \
                            'ctt' in dct_shp['cll'][i_row][i_clm] and \
                            not isnull(dct_shp['cll'][i_row][i_clm]['ctt']) else ''
                    if i_fit not in dct_shp['cll'][i_row][i_clm]:
                        dct_shp['cll'][i_row][i_clm][i_fit] = dct_shp[str_tgt][i_fit][i_row][i_clm]

    def fit_cll(self, dct_shp):
        """
        比较方便地更新self.lcn.ppt.sld.shp.cll
        """
        # 行单元格格式批量处理
        self.fit_mtr(dct_shp, 'cll_fit', dct_shp['clm'])
        self._it_cll(dct_shp, 'cll_fit')
        # 列单元格格式批量修改
        self.fit_mtr(dct_shp, 'cll_fit_clm', dct_shp['row'])
        dct_shp['cll'] = list(map(list, zip(*dct_shp['cll'])))  # 列格式批量处理, 先转置, 再转置还原
        self._it_cll(dct_shp, 'cll_fit_clm')  # 使用dct_shp.cll_fit_clm中的数据进行批量格式赋予
        dct_shp['cll'] = list(map(list, zip(*dct_shp['cll'])))

    def gnr_mrg(self, dct_shp):
        """
        单元格合并, 参数使用dct_shp['mrg']
        """
        dct_shp['mrg'] = self.dct_mtr['tbl']['mrg'] if 'mrg' not in dct_shp else dct_shp['mrg']  # meaningless
        for i_mrg in dct_shp['mrg']:
            dct_shp['obj_tbl'].cell(i_mrg[0], i_mrg[1]).merge(dct_shp['obj_tbl'].cell(i_mrg[2], i_mrg[3]))

    def gnr_cll(self, dct_shp):
        """"""
        dct_shp['obj_tbl'].first_row = False  # 第一行特殊化
        dct_shp['obj_tbl'].first_col = False  # 第一列特殊化
        lst_cll = lsz(dct_shp['cll']).nfd(rtn=True)  # 单元格展开成一行方便后续的赋值操作
        n = 0
        for i_cll in dct_shp['obj_tbl'].iter_cells():
            lst_cll[n] = {'ctt': lst_cll[n]} if type(lst_cll[n]) not in [dict] else lst_cll[n]
            lst_cll[n]['ctt'] = str(lst_cll[n]['ctt']) if 'ctt' in lst_cll[n] and not isnull(lst_cll[n]['ctt']) else ''
            bkg_rgb = lst_cll[n]['color'] if 'color' in lst_cll[n] else dct_shp['color']
            i_cll.fill.solid()
            i_cll.fill.fore_color.rgb = RGBColor(bkg_rgb[0], bkg_rgb[1], bkg_rgb[2])  # 单元格背景色
            i_cll.margin_left = self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['margin_left'])
            i_cll.margin_right = self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['margin_right'])
            i_cll.vertical_anchor = self.dct_anchor[dct_shp['anchor']]  # 纵向对齐
            prh = i_cll.text_frame.paragraphs[0]                # 进入单元格文本处理
            prh.alignment = self.dct_align[dct_shp['align']]    # 横向对齐
            prh.font.size = Pt(lst_cll[n]['fontSize']) if 'fontSize' in lst_cll[n] else Pt(dct_shp['fontSize'])
            prh.font.name = lst_cll[n]['fontName'] if 'fontName' in lst_cll[n] else dct_shp['fontName']
            prh.font.bold = lst_cll[n]['fontBold'] if 'fontBold' in lst_cll[n] else dct_shp['fontBold']
            prh.font.italic = lst_cll[n]['fontItalic'] if 'fontItalic' in lst_cll[n] else dct_shp['fontItalic']
            fnt_rgb = lst_cll[n]['fontColor'] if 'fontColor' in lst_cll[n] else dct_shp['fontColor']
            prh.font.color.rgb = RGBColor(fnt_rgb[0], fnt_rgb[1], fnt_rgb[2])
            prh.text = lst_cll[n]['ctt'] if 'ctt' in lst_cll[n] else ''
            n += 1

    def gnr_tbl(self, ndx_sld, dct_shp):
        """"""
        self.tbl_mtr(dct_shp)
        dct_shp['obj_tbl'] = self.lcn['ppt'][ndx_sld]['obj_shp'].add_table(
            dct_shp['row'],
            dct_shp['clm'],
            self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['left']),
            self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['top']),
            self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['width']),
            self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['height']),
        ).table
        self.fit_cll(dct_shp)
        self.gnr_cll(dct_shp)
        self.gnr_mrg(dct_shp)


class pltMixin(ppz_nit):

    def __init__(self, dts=None, lcn=None, *, spr=False, cvr=True):
        super(pltMixin, self).__init__(dts, lcn, spr=spr, cvr=cvr)

    def gnr_ttl(self, dct_shp):
        """
        generate title, needs 'ttl' in dct_shp
        """
        dct_m = self.lcn['ppt']['mtr']
        if 'ttl' in dct_shp.keys():
            dct_shp['obj_plt'].has_title = True  # 图表是否含有标题，默认为False
            dct_shp['obj_plt'].chart_title.text_frame.clear()  # 清除原标题
            prh = dct_shp['obj_plt'].chart_title.text_frame.paragraphs[0]  # 添加一行新标题
            run = prh.add_run()
            run.text = dct_shp['ttl']
            run.font.size = Pt(dct_shp['titleFontSize']) if 'titleFontSize' in dct_shp else dct_m['fontSize']
            run.font.color.rgb = RGBColor(*self.dct_plt['gry_dep'])
            run.font.name = dct_m['fontName']
            run.font.bold = dct_m['fontBold']
        else:
            dct_shp['obj_plt'].has_title = False

    def gnr_lgd(self, dct_shp):
        """
        generate legend, needs 'lgd' in dct_shp, ['rightout','rightin','leftout/in','bottomout/in','topout/in']
        """
        dct_m = self.lcn['ppt']['mtr']
        if 'lgd' in dct_shp.keys():
            dct_shp['obj_plt'].has_legend = True
            dct_shp['obj_plt'].legend.position = self.dct_legend_position[dct_shp['lgd'].lower()][0]
            dct_shp['obj_plt'].legend.include_in_layout = self.dct_legend_position[dct_shp['lgd']][1]
            dct_shp['obj_plt'].legend.font.name = dct_m['fontName']
            dct_shp['obj_plt'].legend.font.size = dct_m['fontSize_label']
            dct_shp['obj_plt'].legend.font.bold = False
            dct_shp['obj_plt'].legend.font.italic = False
            dct_shp['obj_plt'].legend.font.color.rgb = RGBColor(*self.dct_plt['gry_dep'])
        else:
            dct_shp['obj_plt'].has_legend = False

    def gnr_axs(self, dct_shp):
        """
        generate axis in category and value
        """
        dct_m = self.lcn['ppt']['mtr']
        for i_axs in [dct_shp['obj_plt'].category_axis, dct_shp['obj_plt'].value_axis]:
            i_axs.has_major_gridlines = False
            i_axs.has_minor_gridlines = False
            i_axs.minor_tick_mark = XL_TICK_MARK.NONE
            i_axs.major_tick_mark = XL_TICK_MARK.NONE
            i_axs.tick_labels.font.name = dct_m['fontName']
            i_axs.tick_labels.font.size = dct_m['fontSize_label'] if \
                i_axs in [dct_shp['obj_plt'].category_axis] else dct_m['fontSize_label_2']
            i_axs.tick_labels.font.bold = False
            i_axs.tick_labels.font.italic = False
            i_axs.tick_labels.font.color.rgb = RGBColor(*self.dct_plt['gry_dep'])
            i_axs.format.line.color.rgb = RGBColor(*self.dct_plt['gry'])
            # i_axs.maximum_scale = 100.0  # 纵坐标最大值
            # i_axs.minimum_scale = 0.0  # 纵坐标最小值

    def gnr_lbl(self, dct_shp):
        """
        generate label in a plot
        bar charts in ['outside_end', 'inside_base'] for [柱外顶端, 柱内贴进横轴]
        """
        dct_m = self.lcn['ppt']['mtr']
        plt = dct_shp['obj_plt'].plots[0]
        if 'lbl' in dct_shp:
            plt.has_data_labels = True
            plt.data_labels.font.name = dct_m['fontName']
            plt.data_labels.font.size = dct_m['fontSize_label']
            plt.data_labels.font.color.rgb = RGBColor(*self.dct_plt['gry_dep'])
            plt.data_labels.position = self.dct_label_position[dct_shp['lbl'].lower()]
        else:
            plt.has_data_labels = False

    def gnr_plt(self, ndx_sld, dct_shp):
        chart_data = CategoryChartData()
        chart_data.categories = dct_shp['clm']
        for i in dct_shp['cll']:
            chart_data.add_series(i[0], tuple(i[1:]))
        dct_shp['obj_plt'] = self.lcn['ppt'][ndx_sld]['obj_shp'].add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['left']),
            self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['top']),
            self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['width']),
            self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['height']),
            chart_data,
        ).chart
        dct_shp['obj_plt'].chart_style = dct_shp['sty']  # 图表整体颜色风格[1:48]
        self.gnr_ttl(dct_shp)   # title when dct_shp['ttl']
        self.gnr_axs(dct_shp)   # axis format
        self.gnr_lgd(dct_shp)   # legend
        self.gnr_lbl(dct_shp)   # label


class ppz(tblMixin, pltMixin):
    """
    pptx generator
    >>> lcn = {
    >>>     'fld': '.',
    >>>     'fls': 'tst.pptx',
    >>>     'ppt': {
    >>>         'sld_middle': {
    >>>             'mtr': 'middle',
    >>>             'shp': [
    >>>                 {
    >>>                     'typ': 'textbox',
    >>>                     'tip': '顶栏总结',
    >>>                     'paragraph': [{'ctt': '这里是测试标题', 'fontSize': 16, 'fontBold': True},],
    >>>                 },
    >>>                 {
    >>>                     'typ': 'textbox',
    >>>                     'tip': '左侧内容',
    >>>                     'paragraph': [
    >>>                         {'ctt': '这里是测试的左侧文本内容标题', 'fontSize': 14},
    >>>                         {'ctt': '这里是测试正文，一定要长，测试换行的重要性就在这里体现出来了吗？不知道够不够长了啊，应该不够吧；现在够了吗'},
    >>>                     ],
    >>>                 },
    >>>                 {
    >>>                     'typ': 'table',
    >>>                     'tip': '右侧表格',
    >>>                     'cll': [[12,None,56],[12,34,None],[12,34,56],[{'ctt':12,'color':(225,0,0)},34,56]],
    >>>                     'cll_fit': {
    >>>                         'color': [rgb_rdd,rgb_rdd],
    >>>                         'fontColor': [rgb_wht,rgb_wht],
    >>>                         'fontBold': [True,True],
    >>>                     },
    >>>                     'cll_fit_clm': {'color': [rgb_rdl],},
    >>>                 },
    >>>                 {
    >>>                     'typ': 'picture',
    >>>                     'tip': '右上角logo',
    >>>                     'fld': './pyzohar/samples',
    >>>                     'fls': 'logo_centaline.png',
    >>>                 },
    >>>         ],
    >>>             'tip': '左右平分',
    >>>         },
    >>>         'sld_full': {
    >>>             'mtr': 'full',
    >>>             'shp': [
    >>>                 {
    >>>                     'typ': 'textbox',
    >>>                     'tip': '顶栏总结',
    >>>                     'paragraph': [{'ctt': '这里是测试标题', 'fontSize': 16, 'fontBold': True},],
    >>>                 },
    >>>                 {
    >>>                     'typ':'barchart',
    >>>                     'tip': '主体图表',
    >>>                     'ttl':'演示图表',
    >>>                     'sty': 10,  # [1-48], usually in [9-16], [灰, 彩, 蓝, 红, 绿, 紫, 靛, 橙]
    >>>                     'lbl': 'outside_end',   # 指定数据标签位置, in ['outside_end', 'inside_base']
    >>>                     'lgd': 'rightout',      # 指定图例位置, in ['rightout', 'rightin', 'topout/in', 'bottomout/in']
    >>>                     'clm': ['21-01','21-02','21-03',],  # 横轴分类刻度
    >>>                     'cll': [['ss1',1,2,3],['ss2',1,2,3],],  # 类别和数据
    >>>                 },
    >>>                 {
    >>>                     'typ': 'picture',
    >>>                     'tip': '右上角logo',
    >>>                     'fld': './pyzohar/samples',
    >>>                     'fls': 'logo_centaline.png',
    >>>                 },
    >>>             ],
    >>>             'tip': '全页显示',
    >>>         },
    >>>     },
    >>> }  # sample of mother template 'middle', textbox and table inside
    """
    def __init__(self, dts=None, lcn=None, *, spr=False, cvr=True):
        super(ppz, self).__init__(dts, lcn, spr=spr, cvr=cvr)

    def gnr_pic(self, ndx_sld, dct_shp):
        """
        generate picture
        :param ndx_sld: index of the target slide in self.lcn.ppt
        :param dct_shp: in {'fld','fls','left','top','height','width'}
        """
        w = self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['width']) if 'width' in dct_shp else None
        h = self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['height']) if 'height' in dct_shp else None
        self.lcn['ppt'][ndx_sld]['obj_shp'].add_picture(
            os_join(dct_shp['fld'], dct_shp['fls']),
            self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['left']),
            self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['top']),
            width=w, height=h,
        )

    def gnr_txt(self, ndx_sld, dct_shp):
        """
        generate textbox
        """
        dct_m = self.lcn['ppt']['mtr']
        dct_shp['obj_txt'] = self.lcn['ppt'][ndx_sld]['obj_shp'].add_textbox(
            self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['left']),
            self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['top']),
            self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['width']),
            self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['height']),
        )
        txt = dct_shp['obj_txt'].text_frame
        txt.clear()
        txt.auto_size = MSO_AUTO_SIZE.NONE  # 文本框可以自由拖动，且文字大小与文本框不关联
        txt.vertical_anchor = self.dct_anchor[dct_shp['anchor']] if 'anchor' in dct_shp else self.dct_anchor[
            dct_m['anchor']]  # 纵向对齐（纵向居中）
        txt.word_wrap = True  # 自动换行
        for i_prh in range(len(dct_shp['paragraph'])):
            dct_p = dct_shp['paragraph'][i_prh]
            if i_prh == 0:
                prh = txt.paragraphs[i_prh]
            else:
                prh = txt.add_paragraph()
            prh.line_spacing = self.dct_mtr['line_spacing']  # 行距
            prh.space_before = self.dct_mtr['space_before']  # 段前行距
            prh.space_after = self.dct_mtr['space_after']  # 段后行距
            prh.alignment = self.dct_align[dct_shp['align']] if 'align' in dct_shp else self.dct_align[
                dct_m['align']]  # 文本对齐（两端对齐）
            run = prh.add_run()
            run.text = dct_p['ctt']
            run.font.name = dct_p['fontName'] if 'fontName' in dct_p.keys() else dct_m['fontName']
            fnt_rgb = dct_p['fontColor'] if 'fontColor' in dct_p.keys() else dct_m['fontColor']
            run.font.color.rgb = RGBColor(*fnt_rgb)
            run.font.size = Pt(dct_p['fontSize']) if 'fontSize' in dct_p.keys() else dct_m['fontSize']
            run.font.bold = dct_p['fontBold'] if 'fontBold' in dct_p.keys() else dct_m['fontItalic']
            run.font.italic = dct_p['fontItalic'] if 'fontItalic' in dct_p.keys() else dct_m['fontItalic']

    def gnr_lne(self, ndx_sld, dct_shp):
        """
        generate line
        """
        lns = self.lcn['ppt'][ndx_sld]['obj_shp'].add_shape(
            MSO_SHAPE.RECTANGLE,
            self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['left']),
            self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['top']),
            self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['width']),
            self.dct_util[self.lcn['ppt']['mtr']['utl']](dct_shp['height']),
        )
        lns.line.color.rgb = RGBColor(dct_shp['color'][0], dct_shp['color'][1], dct_shp['color'][2])
        lns.line.color.brightness = 0  # [0,1] => [dark, light]
        lns.line.width = Cm(0.001)
        lns.fill.solid()
        lns.fill.fore_color.rgb = RGBColor(dct_shp['color'][0], dct_shp['color'][1], dct_shp['color'][2])
        lns.shadow.show = False
        lns.shadow.inherit = False

    def gnr_shp(self, ndx_sld):
        """
        generate shapes
        """
        self.lcn['ppt'][ndx_sld]['obj_shp'] = self.lcn['ppt'][ndx_sld]['obj_sld'].shapes
        for dct_k in self.lcn['ppt'][ndx_sld]['shp']:  # 遍历dict of shapes
            if dct_k['typ'] in ['textbox']:
                self.gnr_txt(ndx_sld, dct_k)
            elif dct_k['typ'] in ['line']:
                self.gnr_lne(ndx_sld, dct_k)
            elif dct_k['typ'] in ['picture']:
                self.gnr_pic(ndx_sld, dct_k)
            elif dct_k['typ'] in ['table']:
                self.gnr_tbl(ndx_sld, dct_k)
            elif dct_k['typ'].lower() in ['plot', 'chart', 'barchart', 'linechart', 'piechart']:
                self.gnr_plt(ndx_sld, dct_k)

    def gnr_sld(self):
        """
        generate slides
        """
        for ndx_sld in [i for i in self.lcn['ppt'] if i != 'mtr']:  # ['sld_1']
            print('info: generate slide %s' % self.lcn['ppt'][ndx_sld]['tip'])
            self.sld_mtr(ndx_sld)  # 自定义布局母版检查
            sld_lyt = self.ppt.slide_layouts[6]
            self.lcn['ppt'][ndx_sld]['obj_sld'] = self.ppt.slides.add_slide(sld_lyt)
            self.gnr_shp(ndx_sld)

    def ppt_xpt(self):
        """"""
        self.ppt.save(os_join(self.lcn['fld'], self.lcn['fls']))
